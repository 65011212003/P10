"""Module for exporting presentations to different formats."""

import os
import subprocess
from pathlib import Path
from typing import Optional


class PresentationExporter:
    """Export PowerPoint presentations to various formats."""
    
    @staticmethod
    def export_to_pdf(pptx_path: str, pdf_path: Optional[str] = None) -> str:
        """
        Export PPTX to PDF format.
        
        Args:
            pptx_path: Path to the PPTX file
            pdf_path: Optional output PDF path (defaults to same name as PPTX)
            
        Returns:
            Path to the generated PDF file
            
        Note:
            Requires LibreOffice or Microsoft PowerPoint to be installed
        """
        if pdf_path is None:
            pdf_path = str(Path(pptx_path).with_suffix('.pdf'))
        
        # Try LibreOffice first (cross-platform)
        try:
            subprocess.run([
                'soffice',
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', str(Path(pdf_path).parent),
                pptx_path
            ], check=True, capture_output=True)
            
            # LibreOffice outputs to same directory with .pdf extension
            generated_pdf = str(Path(pptx_path).with_suffix('.pdf'))
            if generated_pdf != pdf_path and Path(generated_pdf).exists():
                os.rename(generated_pdf, pdf_path)
            
            return pdf_path
        except (subprocess.CalledProcessError, FileNotFoundError):
            pass
        
        # Try PowerPoint on Windows
        if os.name == 'nt':
            try:
                import comtypes.client
                
                powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
                powerpoint.Visible = 1
                
                pptx_abs_path = str(Path(pptx_path).resolve())
                pdf_abs_path = str(Path(pdf_path).resolve())
                
                presentation = powerpoint.Presentations.Open(pptx_abs_path)
                presentation.SaveAs(pdf_abs_path, 32)  # 32 = PDF format
                presentation.Close()
                powerpoint.Quit()
                
                return pdf_path
            except Exception:
                pass
        
        raise RuntimeError(
            "PDF export requires LibreOffice (soffice) or Microsoft PowerPoint. "
            "Install LibreOffice: https://www.libreoffice.org/download/"
        )
    
    @staticmethod
    def export_to_html(pptx_path: str, html_path: Optional[str] = None) -> str:
        """
        Export PPTX to HTML5 presentation (reveal.js).
        
        Args:
            pptx_path: Path to the PPTX file
            html_path: Optional output HTML path
            
        Returns:
            Path to the generated HTML file
        """
        if html_path is None:
            html_path = str(Path(pptx_path).with_suffix('.html'))
        
        # Generate basic HTML5 presentation
        from pptx import Presentation
        
        prs = Presentation(pptx_path)
        
        html_content = """<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Presentation</title>
    <link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/reveal.js@4.6.1/dist/reveal.min.css">
    <link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/reveal.js@4.6.1/dist/theme/white.css">
    <style>
        .reveal h1 { font-size: 2.5em; margin-bottom: 0.5em; }
        .reveal h2 { font-size: 2em; margin-bottom: 0.5em; }
        .reveal ul { text-align: left; }
        .reveal li { margin: 0.5em 0; }
    </style>
</head>
<body>
    <div class="reveal">
        <div class="slides">
"""
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            html_content += '            <section>\n'
            
            # Extract text from slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    
                    # First text element as title
                    if slide_num == 1 or '\n' not in text:
                        html_content += f'                <h2>{text}</h2>\n'
                    else:
                        # Multi-line text as list
                        html_content += '                <ul>\n'
                        for line in text.split('\n'):
                            if line.strip():
                                html_content += f'                    <li>{line.strip()}</li>\n'
                        html_content += '                </ul>\n'
            
            html_content += '            </section>\n'
        
        html_content += """        </div>
    </div>
    <script src="https://cdn.jsdelivr.net/npm/reveal.js@4.6.1/dist/reveal.min.js"></script>
    <script>
        Reveal.initialize({
            hash: true,
            transition: 'slide',
            controls: true,
            progress: true,
            center: true,
        });
    </script>
</body>
</html>
"""
        
        with open(html_path, 'w', encoding='utf-8') as f:
            f.write(html_content)
        
        return html_path
    
    @staticmethod
    def get_available_formats() -> dict:
        """Get list of available export formats."""
        return {
            'pdf': 'PDF document (requires LibreOffice or PowerPoint)',
            'html': 'HTML5 presentation (reveal.js)',
        }


def export_presentation(
    pptx_path: str,
    format: str,
    output_path: Optional[str] = None
) -> str:
    """
    Export presentation to specified format.
    
    Args:
        pptx_path: Path to the PPTX file
        format: Export format ('pdf' or 'html')
        output_path: Optional output path
        
    Returns:
        Path to exported file
    """
    exporter = PresentationExporter()
    
    if format == 'pdf':
        return exporter.export_to_pdf(pptx_path, output_path)
    elif format == 'html':
        return exporter.export_to_html(pptx_path, output_path)
    else:
        raise ValueError(f"Unsupported export format: {format}")

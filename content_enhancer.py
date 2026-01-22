"""Enhanced content generation features for presentations."""

from typing import List, Dict, Any, Optional
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pygments import highlight
from pygments.lexers import get_lexer_by_name, guess_lexer
from pygments.formatters import HtmlFormatter
from pygments.util import ClassNotFound


def create_code_slide(
    prs,
    title: str,
    code: str,
    language: str = "python",
    theme: Dict[str, Any] = None
) -> None:
    """
    Create a slide with syntax-highlighted code.
    
    Args:
        prs: PowerPoint presentation object
        title: Slide title
        code: Code content
        language: Programming language for syntax highlighting
        theme: Theme configuration
    """
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    colors = theme['colors']
    
    # Add title
    title_shape = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8)
    )
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.name = theme['font_title']
    title_para.font.size = Pt(36)
    title_para.font.color.rgb = RGBColor(*colors['primary'])
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.LEFT
    
    # Add code box
    code_shape = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.5)
    )
    code_frame = code_shape.text_frame
    code_frame.word_wrap = True
    code_frame.auto_size = MSO_AUTO_SIZE.NONE
    
    # Set monospace font and format
    code_para = code_frame.paragraphs[0]
    code_para.text = code
    code_para.font.name = "Consolas"
    code_para.font.size = Pt(14)
    code_para.font.color.rgb = RGBColor(40, 40, 40)
    
    # Add background color
    fill = code_shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)


def create_table_slide(
    prs,
    title: str,
    headers: List[str],
    rows: List[List[str]],
    theme: Dict[str, Any] = None,
    max_rows: int = 12
) -> None:
    """
    Create a slide with a data table.
    
    Args:
        prs: PowerPoint presentation object
        title: Slide title
        headers: Column headers
        rows: Table data rows
        theme: Theme configuration
        max_rows: Maximum rows to display
    """
    from pptx.oxml.xmlchemy import OxmlElement
    
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    colors = theme['colors']
    
    # Add title
    title_shape = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8)
    )
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.name = theme['font_title']
    title_para.font.size = Pt(36)
    title_para.font.color.rgb = RGBColor(*colors['primary'])
    title_para.font.bold = True
    
    # Limit rows
    display_rows = rows[:max_rows]
    
    # Calculate dimensions
    num_cols = len(headers)
    num_rows = len(display_rows) + 1  # +1 for header
    
    # Add table
    table_shape = slide.shapes.add_table(
        rows=num_rows,
        cols=num_cols,
        left=Inches(0.5),
        top=Inches(1.5),
        width=Inches(12.3),
        height=Inches(5.0)
    )
    
    table = table_shape.table
    
    # Set header row
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*colors['primary'])
        
        # Header text formatting
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(14)
        para.font.bold = True
        para.font.color.rgb = RGBColor(255, 255, 255)
        para.alignment = PP_ALIGN.CENTER
    
    # Set data rows
    for row_idx, row_data in enumerate(display_rows, start=1):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_data)
            
            # Alternate row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
            
            # Cell text formatting
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.font.color.rgb = RGBColor(*colors['text'])
            para.alignment = PP_ALIGN.LEFT


def create_agenda_slide(
    prs,
    sections: List[str],
    theme: Dict[str, Any] = None
) -> None:
    """
    Create an agenda/table of contents slide.
    
    Args:
        prs: PowerPoint presentation object
        sections: List of section titles
        theme: Theme configuration
    """
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    colors = theme['colors']
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = "Agenda"
    for para in title_shape.text_frame.paragraphs:
        para.font.name = theme['font_title']
        para.font.size = Pt(44)
        para.font.color.rgb = RGBColor(*colors['primary'])
        para.font.bold = True
        para.alignment = PP_ALIGN.CENTER
    
    # Set content
    content_shape = slide.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.word_wrap = True
    
    for i, section in enumerate(sections, 1):
        if i == 1:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = f"{i}. {section}"
        p.font.name = theme['font_body']
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(*colors['text'])
        p.space_before = Pt(12)
        p.space_after = Pt(12)


def parse_code_from_content(content: str, language: str = "python") -> Optional[str]:
    """
    Extract code blocks from markdown-style content.
    
    Args:
        content: Content string potentially containing code blocks
        language: Expected programming language
        
    Returns:
        Extracted code or None
    """
    import re
    
    # Look for markdown code blocks
    pattern = r'```(?:' + language + r')?\s*\n(.*?)\n```'
    matches = re.findall(pattern, content, re.DOTALL)
    
    if matches:
        return matches[0]
    
    # Look for indented code blocks
    lines = content.split('\n')
    code_lines = []
    in_code = False
    
    for line in lines:
        if line.startswith('    ') or line.startswith('\t'):
            code_lines.append(line[4:] if line.startswith('    ') else line[1:])
            in_code = True
        elif in_code and not line.strip():
            code_lines.append('')
        elif in_code:
            break
    
    if code_lines:
        return '\n'.join(code_lines)
    
    return None


def parse_csv_to_table(csv_content: str, max_rows: int = 12) -> Optional[tuple]:
    """
    Parse CSV content to table data.
    
    Args:
        csv_content: CSV file content
        max_rows: Maximum rows to include
        
    Returns:
        Tuple of (headers, rows) or None
    """
    import csv
    import io
    
    try:
        reader = csv.reader(io.StringIO(csv_content))
        rows = list(reader)
        
        if len(rows) < 2:
            return None
        
        headers = rows[0]
        data_rows = rows[1:max_rows + 1]
        
        return headers, data_rows
    except Exception:
        return None


def detect_content_type(slide_data: Dict[str, Any]) -> str:
    """
    Detect special content type from slide data.
    
    Args:
        slide_data: Slide data dictionary
        
    Returns:
        Content type: 'code', 'table', 'normal'
    """
    title = slide_data.get('title', '').lower()
    content = slide_data.get('content', [])
    
    # Check title for code keywords
    if any(keyword in title for keyword in ['code', 'implementation', 'example', 'snippet', 'function']):
        # Check if content looks like code
        if content and isinstance(content, list):
            first_item = content[0] if content else ""
            if any(char in first_item for char in ['{', '}', '(', ')', ';', '=']):
                return 'code'
    
    # Check title for table keywords
    if any(keyword in title for keyword in ['table', 'data', 'comparison', 'results', 'statistics']):
        return 'table'
    
    return 'normal'

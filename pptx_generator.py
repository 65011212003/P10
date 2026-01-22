"""Module for generating professional PowerPoint presentations."""

from typing import Dict, List, Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from config import (
    THEMES, 
    DEFAULT_SLIDE_WIDTH, 
    DEFAULT_SLIDE_HEIGHT,
    APP_NAME
)


def rgb_color(rgb_tuple: tuple) -> RGBColor:
    """Convert RGB tuple to RGBColor object."""
    return RGBColor(rgb_tuple[0], rgb_tuple[1], rgb_tuple[2])


def apply_text_formatting(
    paragraph, 
    font_name: str, 
    font_size: int, 
    color: tuple,
    bold: bool = False
) -> None:
    """Apply consistent text formatting to a paragraph."""
    paragraph.font.name = font_name
    paragraph.font.size = Pt(font_size)
    paragraph.font.color.rgb = rgb_color(color)
    paragraph.font.bold = bold


def create_title_slide(
    prs: Presentation, 
    title: str, 
    subtitle: str = "",
    theme: Dict[str, Any] = None
) -> None:
    """Create a professional title slide."""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    colors = theme['colors']
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    
    for paragraph in title_shape.text_frame.paragraphs:
        apply_text_formatting(
            paragraph,
            theme['font_title'],
            theme['title_size'],
            colors['primary'],
            bold=True
        )
        paragraph.alignment = PP_ALIGN.CENTER
    
    # Set subtitle
    if len(slide.placeholders) > 1:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
        
        for paragraph in subtitle_shape.text_frame.paragraphs:
            apply_text_formatting(
                paragraph,
                theme['font_body'],
                theme['subtitle_size'],
                colors['secondary']
            )
            paragraph.alignment = PP_ALIGN.CENTER


def create_content_slide(
    prs: Presentation, 
    title: str, 
    bullet_points: List[str], 
    notes: str = "",
    theme: Dict[str, Any] = None
) -> None:
    """Create a content slide with styled bullet points."""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    colors = theme['colors']
    
    # Set title with styling
    title_shape = slide.shapes.title
    title_shape.text = title
    
    for paragraph in title_shape.text_frame.paragraphs:
        apply_text_formatting(
            paragraph,
            theme['font_title'],
            36,  # Slightly smaller for content slides
            colors['primary'],
            bold=True
        )
    
    # Set content with bullet points
    content_shape = slide.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        # Handle nested points (prefixed with "-" or "*" or indented)
        level = 0
        clean_point = point.strip()
        
        # Detect sub-bullets
        while clean_point.startswith(('  ', '\t', '- ', '• ', '* ')):
            if clean_point.startswith(('  ', '\t')):
                level = min(level + 1, 2)
                clean_point = clean_point[2:] if clean_point.startswith('  ') else clean_point[1:]
            else:
                clean_point = clean_point[2:]
        
        p.text = clean_point
        p.level = level
        
        # Apply formatting based on level
        font_size = theme['bullet_size'] - (level * 2)
        apply_text_formatting(
            p,
            theme['font_body'],
            font_size,
            colors['text']
        )
        p.space_before = Pt(6)
        p.space_after = Pt(4)
    
    # Add speaker notes
    if notes:
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text_frame.text = notes


def create_section_slide(
    prs: Presentation, 
    title: str,
    subtitle: str = "",
    theme: Dict[str, Any] = None
) -> None:
    """Create a section header/divider slide."""
    slide_layout = prs.slide_layouts[2]
    slide = prs.slides.add_slide(slide_layout)
    
    colors = theme['colors']
    
    # Set section title
    title_shape = slide.shapes.title
    title_shape.text = title
    
    for paragraph in title_shape.text_frame.paragraphs:
        apply_text_formatting(
            paragraph,
            theme['font_title'],
            theme['title_size'],
            colors['accent'],
            bold=True
        )
        paragraph.alignment = PP_ALIGN.CENTER


def create_two_column_slide(
    prs: Presentation,
    title: str,
    left_content: List[str],
    right_content: List[str],
    left_title: str = "",
    right_title: str = "",
    theme: Dict[str, Any] = None
) -> None:
    """Create a two-column comparison slide."""
    slide_layout = prs.slide_layouts[3]  # Two Content layout
    slide = prs.slides.add_slide(slide_layout)
    
    colors = theme['colors']
    
    # Set main title
    title_shape = slide.shapes.title
    title_shape.text = title
    
    for paragraph in title_shape.text_frame.paragraphs:
        apply_text_formatting(
            paragraph,
            theme['font_title'],
            36,
            colors['primary'],
            bold=True
        )
    
    # Find content placeholders
    content_shapes = [s for s in slide.placeholders if s.placeholder_format.idx > 0]
    
    if len(content_shapes) >= 2:
        # Left column
        left_frame = content_shapes[0].text_frame
        if left_title:
            left_frame.paragraphs[0].text = left_title
            apply_text_formatting(
                left_frame.paragraphs[0],
                theme['font_title'],
                20,
                colors['secondary'],
                bold=True
            )
            for point in left_content:
                p = left_frame.add_paragraph()
                p.text = point
                apply_text_formatting(p, theme['font_body'], theme['bullet_size'], colors['text'])
        else:
            for i, point in enumerate(left_content):
                if i == 0:
                    left_frame.paragraphs[0].text = point
                    apply_text_formatting(
                        left_frame.paragraphs[0],
                        theme['font_body'],
                        theme['bullet_size'],
                        colors['text']
                    )
                else:
                    p = left_frame.add_paragraph()
                    p.text = point
                    apply_text_formatting(p, theme['font_body'], theme['bullet_size'], colors['text'])
        
        # Right column
        right_frame = content_shapes[1].text_frame
        if right_title:
            right_frame.paragraphs[0].text = right_title
            apply_text_formatting(
                right_frame.paragraphs[0],
                theme['font_title'],
                20,
                colors['secondary'],
                bold=True
            )
            for point in right_content:
                p = right_frame.add_paragraph()
                p.text = point
                apply_text_formatting(p, theme['font_body'], theme['bullet_size'], colors['text'])
        else:
            for i, point in enumerate(right_content):
                if i == 0:
                    right_frame.paragraphs[0].text = point
                    apply_text_formatting(
                        right_frame.paragraphs[0],
                        theme['font_body'],
                        theme['bullet_size'],
                        colors['text']
                    )
                else:
                    p = right_frame.add_paragraph()
                    p.text = point
                    apply_text_formatting(p, theme['font_body'], theme['bullet_size'], colors['text'])


def create_thank_you_slide(
    prs: Presentation,
    title: str = "Thank You",
    subtitle: str = "",
    theme: Dict[str, Any] = None
) -> None:
    """Create a closing thank you slide."""
    slide_layout = prs.slide_layouts[0]  # Title layout
    slide = prs.slides.add_slide(slide_layout)
    
    colors = theme['colors']
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    
    for paragraph in title_shape.text_frame.paragraphs:
        apply_text_formatting(
            paragraph,
            theme['font_title'],
            theme['title_size'] + 4,
            colors['accent'],
            bold=True
        )
        paragraph.alignment = PP_ALIGN.CENTER
    
    # Set subtitle
    if len(slide.placeholders) > 1:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle or "Questions & Discussion"
        
        for paragraph in subtitle_shape.text_frame.paragraphs:
            apply_text_formatting(
                paragraph,
                theme['font_body'],
                theme['subtitle_size'],
                colors['secondary']
            )
            paragraph.alignment = PP_ALIGN.CENTER


def validate_presentation_data(data: Dict[str, Any]) -> Dict[str, Any]:
    """Validate and clean presentation data."""
    if not isinstance(data, dict):
        raise ValueError("Presentation data must be a dictionary")
    
    # Ensure required fields
    if 'title' not in data or not data['title']:
        data['title'] = 'Untitled Presentation'
    
    if 'slides' not in data:
        data['slides'] = []
    
    # Validate and clean slides
    cleaned_slides = []
    for slide in data['slides']:
        if not isinstance(slide, dict):
            continue
            
        cleaned_slide = {
            'title': slide.get('title', 'Untitled Slide'),
            'content': slide.get('content', []),
            'notes': slide.get('notes', ''),
            'type': slide.get('type', 'content'),  # content, section, comparison
        }
        
        # Ensure content is a list
        if isinstance(cleaned_slide['content'], str):
            cleaned_slide['content'] = [cleaned_slide['content']]
        
        # Filter out empty content
        cleaned_slide['content'] = [
            c for c in cleaned_slide['content'] 
            if c and str(c).strip()
        ]
        
        if cleaned_slide['content'] or cleaned_slide['type'] == 'section':
            cleaned_slides.append(cleaned_slide)
    
    data['slides'] = cleaned_slides
    return data


def generate_pptx(
    presentation_data: Dict[str, Any], 
    output_path: str,
    theme_name: str = "professional"
) -> str:
    """
    Generate a professional PowerPoint file from the presentation data.
    
    Args:
        presentation_data: Dictionary containing title and slides.
        output_path: Path where the PPTX file should be saved.
        theme_name: Name of the theme to apply.
        
    Returns:
        The path to the generated PPTX file.
    """
    # Get theme configuration
    theme = THEMES.get(theme_name, THEMES['professional'])
    
    # Validate data
    presentation_data = validate_presentation_data(presentation_data)
    
    # Create presentation
    prs = Presentation()
    
    # Set slide dimensions (16:9 widescreen)
    prs.slide_width = Inches(DEFAULT_SLIDE_WIDTH)
    prs.slide_height = Inches(DEFAULT_SLIDE_HEIGHT)
    
    # Get presentation metadata
    main_title = presentation_data.get('title', 'Untitled Presentation')
    slides = presentation_data.get('slides', [])
    
    # Create title slide
    subtitle = f"Generated by {APP_NAME}"
    create_title_slide(prs, main_title, subtitle, theme)
    
    # Create content slides
    for slide_data in slides:
        slide_type = slide_data.get('type', 'content')
        slide_title = slide_data.get('title', 'Untitled Slide')
        content = slide_data.get('content', [])
        notes = slide_data.get('notes', '')
        
        if slide_type == 'section':
            create_section_slide(prs, slide_title, theme=theme)
        elif slide_type == 'comparison' and 'left' in slide_data and 'right' in slide_data:
            create_two_column_slide(
                prs, 
                slide_title,
                slide_data.get('left', []),
                slide_data.get('right', []),
                slide_data.get('left_title', ''),
                slide_data.get('right_title', ''),
                theme=theme
            )
        else:
            create_content_slide(prs, slide_title, content, notes, theme)
    
    # Create thank you slide
    create_thank_you_slide(prs, theme=theme)
    
    # Save the presentation
    prs.save(output_path)
    
    return output_path

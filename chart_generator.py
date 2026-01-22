"""Module for generating charts from data."""

import io
from typing import List, Dict, Any, Optional, Tuple
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData


class ChartGenerator:
    """Generate charts from data for PowerPoint presentations."""
    
    @staticmethod
    def create_bar_chart(
        slide,
        data: Dict[str, List[float]],
        categories: List[str],
        title: str = "",
        left: float = 1.0,
        top: float = 2.0,
        width: float = 8.0,
        height: float = 4.5
    ):
        """
        Create a bar chart on a slide.
        
        Args:
            slide: PowerPoint slide object
            data: Dictionary mapping series names to values
            categories: List of category labels
            title: Chart title
            left, top, width, height: Chart position and size in inches
        """
        chart_data = CategoryChartData()
        chart_data.categories = categories
        
        for series_name, values in data.items():
            chart_data.add_series(series_name, values)
        
        x, y, cx, cy = Inches(left), Inches(top), Inches(width), Inches(height)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart
        
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
        
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        
        return chart
    
    @staticmethod
    def create_line_chart(
        slide,
        data: Dict[str, List[float]],
        categories: List[str],
        title: str = "",
        left: float = 1.0,
        top: float = 2.0,
        width: float = 8.0,
        height: float = 4.5
    ):
        """
        Create a line chart on a slide.
        
        Args:
            slide: PowerPoint slide object
            data: Dictionary mapping series names to values
            categories: List of category labels
            title: Chart title
            left, top, width, height: Chart position and size in inches
        """
        chart_data = CategoryChartData()
        chart_data.categories = categories
        
        for series_name, values in data.items():
            chart_data.add_series(series_name, values)
        
        x, y, cx, cy = Inches(left), Inches(top), Inches(width), Inches(height)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
        ).chart
        
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
        
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        
        return chart
    
    @staticmethod
    def create_pie_chart(
        slide,
        data: Dict[str, float],
        title: str = "",
        left: float = 2.0,
        top: float = 2.0,
        width: float = 6.0,
        height: float = 4.5
    ):
        """
        Create a pie chart on a slide.
        
        Args:
            slide: PowerPoint slide object
            data: Dictionary mapping labels to values
            title: Chart title
            left, top, width, height: Chart position and size in inches
        """
        chart_data = CategoryChartData()
        chart_data.categories = list(data.keys())
        chart_data.add_series('Values', list(data.values()))
        
        x, y, cx, cy = Inches(left), Inches(top), Inches(width), Inches(height)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
        ).chart
        
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
        
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        
        return chart
    
    @staticmethod
    def parse_csv_for_chart(csv_content: str, max_rows: int = 10) -> Optional[Tuple[List[str], Dict[str, List[float]]]]:
        """
        Parse CSV content to extract chart data.
        
        Args:
            csv_content: CSV file content as string
            max_rows: Maximum number of rows to include
            
        Returns:
            Tuple of (categories, data) or None if parsing fails
        """
        import csv
        import io
        
        try:
            reader = csv.reader(io.StringIO(csv_content))
            rows = list(reader)
            
            if len(rows) < 2:
                return None
            
            # First row is headers
            headers = rows[0]
            
            # First column is categories
            categories = []
            data = {header: [] for header in headers[1:]}
            
            for row in rows[1:max_rows + 1]:
                if len(row) < 2:
                    continue
                
                categories.append(row[0])
                
                for i, header in enumerate(headers[1:], start=1):
                    try:
                        value = float(row[i]) if i < len(row) else 0.0
                        data[header].append(value)
                    except (ValueError, IndexError):
                        data[header].append(0.0)
            
            return categories, data
            
        except Exception:
            return None
    
    @staticmethod
    def detect_chart_type(data: Dict[str, List[float]]) -> str:
        """
        Detect appropriate chart type based on data characteristics.
        
        Args:
            data: Dictionary mapping series names to values
            
        Returns:
            Chart type: 'bar', 'line', or 'pie'
        """
        # If single series, prefer pie chart for small datasets
        if len(data) == 1:
            values = list(data.values())[0]
            if len(values) <= 5:
                return 'pie'
        
        # For multiple series or time-based data, use line chart
        if len(data) > 1:
            return 'line'
        
        # Default to bar chart
        return 'bar'


def add_chart_slide(
    presentation,
    chart_data: Dict[str, Any],
    theme: Dict[str, Any]
) -> None:
    """
    Add a chart slide to a presentation.
    
    Args:
        presentation: PowerPoint presentation object
        chart_data: Dictionary containing chart configuration
        theme: Theme configuration
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank layout
    
    # Add title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = chart_data.get('title', 'Chart')
        title_frame = title_shape.text_frame
        title_frame.paragraphs[0].font.size = Pt(theme['title_size'])
        title_frame.paragraphs[0].font.name = theme['font_title']
        title_frame.paragraphs[0].font.color.rgb = RGBColor(*theme['colors']['primary'])
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Create chart
    chart_type = chart_data.get('type', 'bar')
    categories = chart_data.get('categories', [])
    data = chart_data.get('data', {})
    
    generator = ChartGenerator()
    
    if chart_type == 'bar':
        generator.create_bar_chart(slide, data, categories, "")
    elif chart_type == 'line':
        generator.create_line_chart(slide, data, categories, "")
    elif chart_type == 'pie':
        # Convert data format for pie chart
        pie_data = {cat: vals[0] for cat, vals in zip(categories, data.values())}
        generator.create_pie_chart(slide, pie_data, "")
